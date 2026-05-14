"""
飞书文件消息处理插件 v3

功能：当用户通过飞书发送文件（docx、pdf、xlsx 等）时，
      自动下载并提取文本内容，让 LLM 能直接读取文件内容。

v3 改进（相比 v2）：
- 使用 event.set_extra/get_extra 传递文件内容，避免 session_id 不匹配和实例变量问题
- 简化 on_llm_request 注入逻辑，从 event extra 读取而非实例缓存
- 大幅增强日志输出，方便远程排查 handler 是否被触发
- 优先检查 Comp.File 和 Comp.Image（图片也视为文件）
- 处理 Comp.Reply 中嵌套的文件组件
- 更健壮的文件路径获取（get_file → file_ → url fallback）
"""

import os
from pathlib import Path

import astrbot.api.message_components as Comp
from astrbot.api import logger
from astrbot.api.event import AstrMessageEvent, filter
from astrbot.api.provider import ProviderRequest
from astrbot.api.star import Context, Star
from astrbot.core.star.filter.event_message_type import EventMessageType
from astrbot.core.star.filter.platform_adapter_type import PlatformAdapterType


# 默认最大提取文件大小：10MB
MAX_FILE_SIZE = 10 * 1024 * 1024

# 最大单文件提取内容长度（字符数）
MAX_CONTENT_LENGTH = 50000

# Extra key for passing file contents between handlers
_EXTRA_KEY = "_lark_file_handler_contents"

# 支持的文件格式及对应的提取方式
SUPPORTED_FORMATS: dict[str, str] = {
    # 文本类
    ".txt": "text",
    ".md": "text",
    ".csv": "text",
    ".json": "text",
    ".py": "text",
    ".js": "text",
    ".ts": "text",
    ".yaml": "text",
    ".yml": "text",
    ".xml": "text",
    ".html": "text",
    ".css": "text",
    ".sql": "text",
    ".log": "text",
    ".ini": "text",
    ".conf": "text",
    ".env": "text",
    ".toml": "text",
    # 文档类
    ".docx": "docx",
    ".pdf": "pdf",
    ".xlsx": "xlsx",
    ".pptx": "pptx",
}


class LarkFileHandler(Star):
    """飞书文件消息处理器 v3"""

    def __init__(self, context: Context):
        super().__init__(context)
        logger.info("[LarkFile] 插件已初始化，等待飞书文件消息...")

    @filter.platform_adapter_type(PlatformAdapterType.LARK)
    @filter.event_message_type(EventMessageType.ALL)
    async def on_file_message(self, event: AstrMessageEvent):
        """拦截飞书文件消息，自动下载并提取内容。

        通过 event.set_extra 将提取的内容传递给 on_llm_request 钩子，
        同时修改 event.message_str 作为备用注入路径。
        """
        logger.info(
            f"[LarkFile] on_file_message 被触发 | "
            f"platform={event.get_platform_name()} | "
            f"session={event.session_id} | "
            f"msg_str={event.message_str!r[:100]}"
        )

        try:
            # 收集所有文件组件（包括 Reply 中嵌套的）
            file_components = []
            for comp in event.message_obj.message:
                if isinstance(comp, Comp.File):
                    file_components.append(comp)
                elif isinstance(comp, Comp.Image):
                    # 图片也提取（可能有 OCR 需求，暂跳过，只记录）
                    logger.debug(f"[LarkFile] 发现图片组件: {comp}")
                elif isinstance(comp, Comp.Reply) and comp.chain:
                    for reply_comp in comp.chain:
                        if isinstance(reply_comp, Comp.File):
                            file_components.append(reply_comp)

            if not file_components:
                logger.debug("[LarkFile] 消息中无文件组件，跳过")
                return

            logger.info(f"[LarkFile] 发现 {len(file_components)} 个文件组件，开始处理")

            # 处理每个文件
            file_contents = []
            for i, comp in enumerate(file_components):
                logger.info(
                    f"[LarkFile] 处理文件 #{i+1}: name={comp.name!r}, "
                    f"file_={comp.file_!r}, url={comp.url!r}"
                )
                content = await self._process_file(comp)
                if content:
                    file_contents.append(content)

            if not file_contents:
                logger.info("[LarkFile] 未提取到任何文件内容")
                return

            combined_content = "\n\n---\n\n".join(file_contents)

            # 路径1: 通过 event extra 传递（供 on_llm_request 使用）
            existing = event.get_extra(_EXTRA_KEY, [])
            existing.append(combined_content)
            event.set_extra(_EXTRA_KEY, existing)

            # 路径2: 修改 message_str（兼容不走 LLM 的场景）
            original_text = event.message_str or ""
            inject_text = (
                f"\n\n[文件内容已自动提取，以下是文件内容]\n\n{combined_content}"
            )
            event.message_str = original_text + inject_text

            logger.info(
                f"[LarkFile] 文件内容已注入 | "
                f"总长度={len(combined_content)} | "
                f"message_str 长度={len(event.message_str)}"
            )

        except Exception as e:
            logger.error(f"[LarkFile] on_file_message 处理失败: {e}", exc_info=True)

    @filter.on_llm_request()
    async def inject_file_content_to_llm(
        self, event: AstrMessageEvent, req: ProviderRequest
    ):
        """在 LLM 请求前注入文件内容。

        从 event extra 中读取 on_file_message 缓存的内容，
        追加到 ProviderRequest.prompt 中。
        """
        cached = event.get_extra(_EXTRA_KEY, None)
        if not cached:
            return

        # 取出后清除，避免重复注入
        event.set_extra(_EXTRA_KEY, None)

        combined = "\n\n---\n\n".join(cached)

        # 追加到 prompt 末尾
        if req.prompt:
            req.prompt += f"\n\n[文件内容已自动提取]\n\n{combined}"
        else:
            req.prompt = f"[文件内容已自动提取]\n\n{combined}"

        # 如果原始 prompt 为空，补充一个默认提示
        if not req.prompt.strip().replace("[文件内容已自动提取]", "").strip():
            req.prompt += "\n\n请总结文件内容。"

        logger.info(
            f"[LarkFile] 已通过 on_llm_request 注入文件内容 | "
            f"prompt 长度={len(req.prompt)}"
        )

    async def _process_file(self, comp: Comp.File) -> str | None:
        """处理单个文件组件，返回提取的文本内容"""
        try:
            # 获取文件路径：优先异步方法，然后 fallback 到内部字段
            file_path = ""
            try:
                file_path = await comp.get_file()
            except Exception as e:
                logger.warning(f"[LarkFile] get_file() 失败: {e}, 尝试 fallback")

            if not file_path:
                file_path = comp.file_ or ""

            if not file_path and comp.url:
                # 最后的 fallback: 只有 URL，无法本地提取
                logger.warning(
                    f"[LarkFile] 文件仅有 URL，无法本地提取: {comp.url[:80]}"
                )
                return None

            file_name = comp.name or "unknown"

            if not file_path or not os.path.exists(file_path):
                logger.warning(
                    f"[LarkFile] 文件路径不存在: {file_path!r} (name={file_name})"
                )
                return None

            # 文件大小检查
            file_size = os.path.getsize(file_path)
            if file_size > MAX_FILE_SIZE:
                logger.warning(
                    f"[LarkFile] 文件过大 ({file_size / 1024 / 1024:.1f}MB)，跳过: {file_name}"
                )
                return (
                    f"📄 **{file_name}**: 文件过大"
                    f"（{file_size / 1024 / 1024:.1f}MB），无法自动提取内容。"
                )

            # 获取文件扩展名
            ext = Path(file_path).suffix.lower()
            if ext not in SUPPORTED_FORMATS:
                logger.info(f"[LarkFile] 不支持的文件格式: {ext} (file={file_name})")
                return None

            # 根据格式提取内容
            file_type = SUPPORTED_FORMATS[ext]
            content = None

            if file_type == "text":
                content = await self._read_text_file(file_path)
            elif file_type == "docx":
                content = await self._read_docx(file_path)
            elif file_type == "pdf":
                content = await self._read_pdf(file_path)
            elif file_type == "xlsx":
                content = await self._read_excel(file_path)
            elif file_type == "pptx":
                content = await self._read_pptx(file_path)

            if content:
                # 截断超长内容
                if len(content) > MAX_CONTENT_LENGTH:
                    content = content[:MAX_CONTENT_LENGTH] + "\n\n[内容已截断...]"
                logger.info(
                    f"[LarkFile] 提取成功: {file_name} | "
                    f"类型={file_type} | 长度={len(content)}"
                )
                return f"📄 **{file_name}** 的内容：\n\n{content}"

            logger.warning(f"[LarkFile] 提取内容为空: {file_name} (类型={file_type})")
            return None

        except Exception as e:
            logger.error(f"[LarkFile] _process_file 失败: {e}", exc_info=True)
            return None

    async def _read_text_file(self, file_path: str) -> str | None:
        """读取文本文件，自动尝试 UTF-8 / GBK 编码"""
        for encoding in ("utf-8", "gbk", "latin-1"):
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logger.error(f"[LarkFile] 读取文本文件失败: {e}")
                return None
        return None

    async def _read_docx(self, file_path: str) -> str | None:
        """读取 Word 文档"""
        try:
            from docx import Document

            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            # 读取表格
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    paragraphs.append(" | ".join(cells))
            return "\n".join(paragraphs)
        except ImportError:
            logger.warning("[LarkFile] python-docx 未安装，无法读取 docx 文件")
            return None
        except Exception as e:
            logger.error(f"[LarkFile] 读取 docx 失败: {e}")
            return None

    async def _read_pdf(self, file_path: str) -> str | None:
        """读取 PDF 文件，优先 pdfplumber，回退到 PyPDF2。"""
        # 尝试 pdfplumber（提取质量更好）
        try:
            import pdfplumber

            text_parts = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)
            if text_parts:
                return "\n\n".join(text_parts)
        except ImportError:
            logger.debug("[LarkFile] pdfplumber 未安装，尝试 PyPDF2")
        except Exception as e:
            logger.warning(f"[LarkFile] pdfplumber 读取 PDF 失败: {e}")

        # 回退到 PyPDF2（纯 Python，适合 ARM）
        try:
            from PyPDF2 import PdfReader

            reader = PdfReader(file_path)
            text_parts = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
            if text_parts:
                return "\n\n".join(text_parts)
        except ImportError:
            logger.warning("[LarkFile] PyPDF2 也未安装，无法读取 PDF 文件")
        except Exception as e:
            logger.error(f"[LarkFile] PyPDF2 读取 PDF 失败: {e}")

        return None

    async def _read_excel(self, file_path: str) -> str | None:
        """读取 Excel 文件"""
        try:
            import openpyxl

            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            text_parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                text_parts.append(f"--- Sheet: {sheet_name} ---")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(cell) if cell is not None else "" for cell in row]
                    text_parts.append(" | ".join(cells))
            wb.close()
            return "\n".join(text_parts)
        except ImportError:
            logger.warning("[LarkFile] openpyxl 未安装，无法读取 xlsx 文件")
            return None
        except Exception as e:
            logger.error(f"[LarkFile] 读取 xlsx 失败: {e}")
            return None

    async def _read_pptx(self, file_path: str) -> str | None:
        """读取 PPT 文件"""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            text_parts = []
            for i, slide in enumerate(prs.slides, 1):
                text_parts.append(f"--- 第 {i} 页 ---")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                text_parts.append(text)
            return "\n".join(text_parts)
        except ImportError:
            logger.warning("[LarkFile] python-pptx 未安装，无法读取 pptx 文件")
            return None
        except Exception as e:
            logger.error(f"[LarkFile] 读取 pptx 失败: {e}")
            return None


star_cls = LarkFileHandler
